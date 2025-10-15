#!/usr/bin/env python3
"""
Demo and test script for Presentation Toolkit.
Creates sample files and demonstrates functionality.
"""

import os
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter

from font_extractor import FontExtractor
from pdf_converter import PDFToPPTXConverter


def create_sample_pptx(filename="sample_presentation.pptx"):
    """Create a sample PowerPoint file for testing."""
    print(f"Creating sample PowerPoint: {filename}")
    
    prs = Presentation()
    
    # Add a title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Sample Presentation"
    subtitle.text = "Created by Presentation Toolkit Demo"
    
    # Add a content slide
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Sample Content"
    
    tf = body_shape.text_frame
    tf.text = "Font extraction demo"
    
    p = tf.add_paragraph()
    p.text = "Testing font detection"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "This file uses system fonts"
    p.level = 2
    
    # Save the presentation
    prs.save(filename)
    print(f"✓ Created: {filename}")
    return filename


def create_sample_pdf(filename="sample_document.pdf"):
    """Create a sample PDF file for testing."""
    print(f"Creating sample PDF: {filename}")
    
    c = canvas.Canvas(filename, pagesize=letter)
    width, height = letter
    
    # Page 1
    c.setFont("Helvetica-Bold", 24)
    c.drawString(100, height - 100, "Sample PDF Document")
    
    c.setFont("Helvetica", 14)
    c.drawString(100, height - 150, "This PDF will be converted to PowerPoint")
    c.drawString(100, height - 180, "Each page becomes a slide")
    
    c.showPage()
    
    # Page 2
    c.setFont("Helvetica-Bold", 20)
    c.drawString(100, height - 100, "Page 2")
    
    c.setFont("Helvetica", 12)
    c.drawString(100, height - 150, "This is the second page")
    c.drawString(100, height - 180, "It will become the second slide")
    
    c.showPage()
    
    # Page 3
    c.setFont("Helvetica-Bold", 20)
    c.drawString(100, height - 100, "Page 3")
    
    c.setFont("Helvetica", 12)
    c.drawString(100, height - 150, "Final page of the demo PDF")
    
    c.save()
    print(f"✓ Created: {filename}")
    return filename


def demo_font_extraction():
    """Demonstrate font extraction from PowerPoint."""
    print("\n" + "="*60)
    print("DEMO: Font Extraction from PowerPoint")
    print("="*60 + "\n")
    
    # Create sample file
    sample_file = create_sample_pptx()
    
    # Extract fonts
    print(f"\nExtracting fonts from {sample_file}...")
    extractor = FontExtractor(output_dir="demo_extracted_fonts")
    
    try:
        result = extractor.extract_from_pptx(sample_file)
        
        print(f"\n✓ Font extraction complete!")
        print(f"  Embedded fonts found: {len(result['embedded_fonts'])}")
        
        if result['embedded_fonts']:
            print("  Extracted fonts:")
            for font in result['embedded_fonts']:
                print(f"    - {Path(font).name}")
        else:
            print("  (No embedded fonts - this is normal for system font usage)")
        
        print(f"\n  Referenced fonts found: {len(result['referenced_fonts'])}")
        if result['referenced_fonts']:
            print("  Font references in presentation:")
            for font in result['referenced_fonts']:
                print(f"    - {font}")
        
        print(f"\n  Output folder: {result['output_folder']}")
        
    except Exception as e:
        print(f"✗ Error: {e}")
        return False
    
    return True


def demo_pdf_conversion():
    """Demonstrate PDF to PowerPoint conversion."""
    print("\n" + "="*60)
    print("DEMO: PDF to PowerPoint Conversion")
    print("="*60 + "\n")
    
    # Create sample PDF
    sample_pdf = create_sample_pdf()
    
    # Convert to PowerPoint
    print(f"\nConverting {sample_pdf} to PowerPoint...")
    converter = PDFToPPTXConverter(output_dir="demo_converted_pptx")
    
    try:
        output_file = converter.convert(sample_pdf, dpi=150)  # Lower DPI for demo
        converter.cleanup_temp()
        
        print(f"\n✓ Conversion complete!")
        print(f"  Created: {output_file}")
        print(f"  You can now open this file in PowerPoint")
        
    except Exception as e:
        print(f"\n✗ Error: {e}")
        print("\nNote: PDF conversion requires Poppler to be installed.")
        print("  macOS: brew install poppler")
        print("  Ubuntu: sudo apt-get install poppler-utils")
        return False
    
    return True


def cleanup_demo_files():
    """Clean up demo files."""
    print("\n" + "="*60)
    print("Cleanup")
    print("="*60 + "\n")
    
    print("Demo files created:")
    print("  - sample_presentation.pptx")
    print("  - sample_document.pdf")
    print("  - demo_extracted_fonts/")
    print("  - demo_converted_pptx/")
    
    response = input("\nDo you want to delete these demo files? (y/N): ")
    
    if response.lower() == 'y':
        files_to_remove = [
            "sample_presentation.pptx",
            "sample_document.pdf"
        ]
        
        for file in files_to_remove:
            if Path(file).exists():
                Path(file).unlink()
                print(f"  Deleted: {file}")
        
        print("\nNote: Demo output folders kept for your review")
        print("  Delete manually if needed: demo_extracted_fonts/, demo_converted_pptx/")
    else:
        print("Demo files kept for your review")


def main():
    """Run the demo."""
    print("\n" + "="*60)
    print("Presentation Toolkit - Demo & Test")
    print("="*60)
    print("\nThis demo will:")
    print("  1. Create sample presentation and PDF files")
    print("  2. Demonstrate font extraction")
    print("  3. Demonstrate PDF to PowerPoint conversion")
    print("\n")
    
    input("Press Enter to start the demo...")
    
    # Run demos
    font_demo_success = demo_font_extraction()
    pdf_demo_success = demo_pdf_conversion()
    
    # Summary
    print("\n" + "="*60)
    print("Demo Summary")
    print("="*60)
    print(f"\nFont Extraction: {'✓ Success' if font_demo_success else '✗ Failed'}")
    print(f"PDF Conversion: {'✓ Success' if pdf_demo_success else '✗ Failed'}")
    
    if font_demo_success and pdf_demo_success:
        print("\n✓ All demos completed successfully!")
        print("\nThe toolkit is ready to use with your files.")
    elif not pdf_demo_success:
        print("\n⚠ PDF conversion failed - make sure Poppler is installed")
        print("  Font extraction is working and ready to use!")
    
    print("\nSee QUICK_REFERENCE.md for common commands")
    print("See README.md for complete documentation")
    
    # Cleanup option
    cleanup_demo_files()


if __name__ == "__main__":
    try:
        main()
    except KeyboardInterrupt:
        print("\n\nDemo cancelled by user")
        sys.exit(0)
    except Exception as e:
        print(f"\n\nError running demo: {e}")
        sys.exit(1)


