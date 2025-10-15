"""
PDF to PowerPoint converter.
Converts PDF pages to images and creates a PowerPoint presentation.
"""

import os
from pathlib import Path
from typing import List, Optional
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
from PIL import Image


class PDFToPPTXConverter:
    """Convert PDF files to PowerPoint presentations."""
    
    def __init__(self, output_dir: str = "converted_pptx", temp_dir: str = "temp"):
        self.output_dir = Path(output_dir)
        self.temp_dir = Path(temp_dir)
        self.output_dir.mkdir(exist_ok=True, parents=True)
        self.temp_dir.mkdir(exist_ok=True, parents=True)
    
    def convert(self, pdf_path: str, dpi: int = 300, output_name: Optional[str] = None) -> str:
        """
        Convert a PDF file to PowerPoint presentation.
        
        Args:
            pdf_path: Path to the PDF file
            dpi: Resolution for PDF to image conversion (default: 300)
            output_name: Optional custom name for output file
            
        Returns:
            Path to the created PowerPoint file
        """
        pdf_path = Path(pdf_path)
        
        if not pdf_path.exists():
            raise FileNotFoundError(f"PDF file not found: {pdf_path}")
        
        # Determine output filename
        if output_name:
            output_filename = output_name if output_name.endswith('.pptx') else f"{output_name}.pptx"
        else:
            output_filename = f"{pdf_path.stem}.pptx"
        
        output_path = self.output_dir / output_filename
        
        try:
            # Convert PDF pages to images
            print(f"Converting PDF to images (DPI: {dpi})...")
            images = convert_from_path(str(pdf_path), dpi=dpi)
            
            # Create PowerPoint presentation
            print(f"Creating PowerPoint with {len(images)} slides...")
            prs = Presentation()
            
            # Set slide dimensions (16:9 widescreen)
            # Standard 16:9 widescreen dimensions: 10" x 5.625" (or 13.33" x 7.5")
            prs.slide_width = Inches(13.33)  # 16:9 widescreen width
            prs.slide_height = Inches(7.5)   # 16:9 widescreen height
            
            # Process each page/image
            for i, image in enumerate(images, 1):
                print(f"  Processing slide {i}/{len(images)}...")
                
                # Save image temporarily
                temp_image_path = self.temp_dir / f"page_{i}.png"
                image.save(str(temp_image_path), 'PNG')
                
                # Add blank slide
                blank_slide_layout = prs.slide_layouts[6]  # Blank layout
                slide = prs.slides.add_slide(blank_slide_layout)
                
                # Calculate image dimensions to fit slide
                img_width, img_height = image.size
                slide_width = prs.slide_width
                slide_height = prs.slide_height
                
                # Calculate scaling to fit image in slide while maintaining aspect ratio
                width_ratio = slide_width / img_width
                height_ratio = slide_height / img_height
                scale_ratio = min(width_ratio, height_ratio)
                
                # Calculate final dimensions
                final_width = int(img_width * scale_ratio)
                final_height = int(img_height * scale_ratio)
                
                # Calculate position to center the image
                left = (slide_width - final_width) // 2
                top = (slide_height - final_height) // 2
                
                # Add image to slide
                slide.shapes.add_picture(
                    str(temp_image_path),
                    left,
                    top,
                    width=final_width,
                    height=final_height
                )
                
                # Clean up temporary image
                temp_image_path.unlink()
            
            # Save presentation
            prs.save(str(output_path))
            print(f"Successfully created: {output_path}")
            
            return str(output_path)
        
        except Exception as e:
            raise Exception(f"Error converting PDF to PPTX: {e}")
    
    def convert_multiple(self, pdf_paths: List[str], dpi: int = 300) -> List[str]:
        """
        Convert multiple PDF files to PowerPoint presentations.
        
        Args:
            pdf_paths: List of paths to PDF files
            dpi: Resolution for PDF to image conversion
            
        Returns:
            List of paths to created PowerPoint files
        """
        output_files = []
        
        for pdf_path in pdf_paths:
            try:
                output_file = self.convert(pdf_path, dpi=dpi)
                output_files.append(output_file)
            except Exception as e:
                print(f"Error converting {pdf_path}: {e}")
        
        return output_files
    
    def cleanup_temp(self):
        """Clean up temporary files."""
        if self.temp_dir.exists():
            for file in self.temp_dir.glob("*"):
                file.unlink()


def convert_pdf_to_pptx(pdf_path: str, output_dir: str = "converted_pptx", dpi: int = 300) -> str:
    """
    Convenience function to convert a PDF to PowerPoint.
    
    Args:
        pdf_path: Path to the PDF file
        output_dir: Directory to save the PowerPoint file
        dpi: Resolution for conversion (default: 300)
        
    Returns:
        Path to the created PowerPoint file
    """
    converter = PDFToPPTXConverter(output_dir=output_dir)
    result = converter.convert(pdf_path, dpi=dpi)
    converter.cleanup_temp()
    return result

